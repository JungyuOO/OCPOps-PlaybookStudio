from __future__ import annotations

import base64
import json
import sys
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

import requests
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "src"
if str(SRC) not in sys.path:
    sys.path.insert(0, str(SRC))

from play_book_studio.app.intake_api import ingest_customer_pack
from tests.test_customer_pack_read_boundary import (
    _FakeChunkingModel,
    _FakeEmbeddingModel,
    _test_server,
)


def _create_docx(path: Path) -> None:
    document = Document()
    document.add_heading("백업 절차", level=1)
    document.add_paragraph("oc get pods -A")
    document.add_paragraph("확인: Pod 상태를 점검한다.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "항목"
    table.cell(0, 1).text = "값"
    table.cell(1, 0).text = "모드"
    table.cell(1, 1).text = "active"
    document.save(path)


def _create_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "운영 점검"
    slide.placeholders[1].text = "oc get nodes\n확인: Node 상태를 점검한다."
    presentation.save(path)


def _create_messy_pptx(path: Path) -> None:
    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    image_path = path.with_suffix(".png")
    image_path.write_bytes(
        base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO9WcXQAAAAASUVORK5CYII="
        )
    )

    slide = presentation.slides.add_slide(blank)
    table = slide.shapes.add_table(1, 1, Inches(1.2), Inches(1.0), Inches(8.0), Inches(1.0)).table
    table.cell(0, 0).text = "KOMSCO 지급결제플랫폼 아키텍처 설계서 (CICD)"
    date_box = slide.shapes.add_textbox(Inches(3.6), Inches(5.9), Inches(2.0), Inches(0.4))
    date_box.text_frame.text = "2025. 07. 25"
    for paragraph in date_box.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(18)

    slide = presentation.slides.add_slide(blank)
    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.3), Inches(5.8), Inches(0.4))
    title_box.text_frame.text = "CICD 프로세스"
    for paragraph in title_box.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(18)
    code_box = slide.shapes.add_textbox(Inches(7.5), Inches(0.3), Inches(2.0), Inches(0.4))
    code_box.text_frame.text = "DSGN-005-401"
    for paragraph in code_box.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(16)

    labels = (
        ("개발(aka. 검증) 환경", 0.9, 1.0, 2.9),
        ("운영 환경", 5.5, 1.0, 2.0),
        ("CI\n(Tekton)", 0.4, 2.1, 0.9),
        ("CD\n(ArgoCD)", 0.4, 3.4, 0.9),
        ("소스 빌드", 2.3, 2.4, 0.9),
        ("이미지빌드", 3.6, 2.4, 0.9),
        ("배포명세 변경", 5.0, 2.4, 1.1),
        ("Push", 3.8, 3.0, 0.7),
        ("Sync", 5.2, 4.2, 0.8),
        ("Pull", 4.1, 5.0, 0.8),
    )
    for text, left, top, width in labels:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.5))
        box.text_frame.text = text
        for paragraph in box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(10)

    slide = presentation.slides.add_slide(blank)
    slide.shapes.add_picture(str(image_path), Inches(1.5), Inches(1.4), width=Inches(4.5), height=Inches(2.8))

    presentation.save(path)


def _create_xlsx(path: Path) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "정책 매트릭스"
    sheet.append(["항목", "값"])
    sheet.append(["지원", "enabled"])
    sheet.append(["제한", "local_only"])
    workbook.save(path)


class CustomerPackNativeOoxmlSmokeTests(unittest.TestCase):
    def test_native_ooxml_files_ingest_and_render_end_to_end(self) -> None:
        cases = (
            ("docx", _create_docx, "docx_native_structure", "procedure"),
            ("pptx", _create_pptx, "pptx_native_slide_extract", "procedure"),
            ("xlsx", _create_xlsx, "xlsx_native_sheet_extract", "reference"),
        )
        for source_type, create_fixture, expected_backend, expected_role in cases:
            with self.subTest(source_type=source_type):
                with tempfile.TemporaryDirectory() as tmpdir:
                    root = Path(tmpdir)
                    source_path = root / f"sample.{source_type}"
                    create_fixture(source_path)

                    with (
                        patch(
                            "play_book_studio.intake.private_corpus.load_sentence_model",
                            return_value=_FakeEmbeddingModel(),
                        ),
                        patch(
                            "play_book_studio.ingestion.chunking.load_sentence_model",
                            return_value=_FakeChunkingModel(),
                        ),
                    ):
                        result = ingest_customer_pack(
                            root,
                            {
                                "source_type": source_type,
                                "uri": str(source_path),
                                "title": f"{source_type.upper()} 샘플",
                                "approval_state": "approved",
                            },
                        )

                    book = dict(result.get("book") or {})
                    evidence = dict(book.get("customer_pack_evidence") or {})
                    sections = [
                        dict(section)
                        for section in (book.get("sections") or [])
                        if isinstance(section, dict)
                    ]

                    self.assertEqual("normalized", result["status"])
                    self.assertTrue(sections)
                    self.assertEqual("native_ooxml_first", evidence["primary_parse_strategy"])
                    self.assertEqual(expected_backend, evidence["parser_backend"])
                    self.assertTrue(any(str(section.get("semantic_role") or "") == expected_role for section in sections))
                    self.assertTrue(any(section.get("block_kinds") for section in sections))

                    with _test_server(root) as (base_url, _store, _answerer):
                        response = requests.get(
                            f"{base_url}/playbooks/customer-packs/{result['draft_id']}/index.html",
                            timeout=10,
                        )

                    self.assertEqual(200, response.status_code)
                    self.assertIn(f"{source_type.upper()} 샘플", response.text)

    def test_messy_pptx_infers_slide_titles_and_filters_document_code(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            source_path = root / "messy.pptx"
            _create_messy_pptx(source_path)

            with (
                patch(
                    "play_book_studio.intake.private_corpus.load_sentence_model",
                    return_value=_FakeEmbeddingModel(),
                ),
                patch(
                    "play_book_studio.ingestion.chunking.load_sentence_model",
                    return_value=_FakeChunkingModel(),
                ),
            ):
                result = ingest_customer_pack(
                    root,
                    {
                        "source_type": "pptx",
                        "uri": str(source_path),
                        "title": "P 유형 샘플",
                        "approval_state": "approved",
                    },
                )

            book = dict(result.get("book") or {})
            sections = [
                dict(section)
                for section in (book.get("sections") or [])
                if isinstance(section, dict)
            ]
            headings = [str(section.get("heading") or "") for section in sections]

            self.assertEqual("normalized", result["status"])
            self.assertIn("KOMSCO 지급결제플랫폼 아키텍처 설계서 (CICD)", headings)
            self.assertIn("CICD 프로세스", headings)
            self.assertNotIn("Slide 1", headings)
            self.assertNotIn("Slide 2", headings)

            process_section = next(section for section in sections if section.get("heading") == "CICD 프로세스")
            self.assertIn("figure", [str(item) for item in (process_section.get("block_kinds") or [])])
            self.assertIn("CI", str(process_section.get("text") or ""))
            self.assertIn("ArgoCD", str(process_section.get("text") or ""))
            self.assertNotIn("DSGN-005-401", str(process_section.get("text") or ""))
            artifact_bundle = dict(book.get("artifact_bundle") or {})
            manifest = json.loads(Path(str(book["artifact_manifest_path"])).read_text(encoding="utf-8"))
            relations_payload = json.loads(Path(str(manifest["relations_path"])).read_text(encoding="utf-8"))
            figure_assets_payload = json.loads(Path(str(manifest["figure_assets_path"])).read_text(encoding="utf-8"))
            slide_packets_payload = json.loads(Path(str(manifest["slide_packets_path"])).read_text(encoding="utf-8"))
            candidate_relations = dict(relations_payload.get("candidate_relations") or {})
            entity_hubs = dict(relations_payload.get("entity_hubs") or {})
            figure_section_index = dict(relations_payload.get("figure_section_index") or {})
            section_relation_index = dict(relations_payload.get("section_relation_index") or {})

            self.assertGreater(int(artifact_bundle.get("relation_count") or 0), 0)
            self.assertGreater(int(manifest.get("relation_count") or 0), len(sections))
            self.assertTrue(candidate_relations)
            self.assertTrue(entity_hubs)
            self.assertNotIn("2025-07-25", entity_hubs)
            self.assertIn("by_book", section_relation_index)
            self.assertTrue(section_relation_index["by_book"][str(book["book_slug"])])
            self.assertIn("by_slug", figure_section_index)
            self.assertTrue(figure_section_index["by_slug"][str(book["book_slug"])])
            self.assertIn(str(book["book_slug"]), figure_assets_payload["entries"])
            self.assertTrue(figure_assets_payload["entries"][str(book["book_slug"])])
            self.assertTrue(any("tekton" in relation_id for relation_id in candidate_relations))
            self.assertEqual("slide_deck", book["surface_kind"])
            self.assertEqual("slide", book["source_unit_kind"])
            self.assertEqual(3, int(book["source_unit_count"]))
            self.assertEqual(3, int(book["slide_packet_count"]))
            self.assertGreaterEqual(int(book["slide_asset_count"]), 1)
            self.assertEqual("slide_deck", slide_packets_payload["surface_kind"])
            self.assertEqual(3, int(slide_packets_payload["slide_count"]))
            self.assertEqual(3, len(slide_packets_payload["slides"]))
            self.assertTrue(any(asset.get("asset_kind") == "image" for asset in slide_packets_payload["embedded_assets"]))
            image_asset = next(asset for asset in slide_packets_payload["embedded_assets"] if asset.get("asset_kind") == "image")
            self.assertTrue((Path(str(manifest["slide_packets_path"])).parent / str(image_asset["storage_relpath"])).exists())
            self.assertEqual("visual_only", slide_packets_payload["slides"][2]["slide_role"])
            self.assertEqual("Slide 3", slide_packets_payload["slides"][2]["title"])
            self.assertFalse(slide_packets_payload["slides"][2]["matched_section_anchor"])

            with _test_server(root) as (base_url, _store, _answerer):
                asset_url = f"/playbooks/customer-packs/{result['draft_id']}/artifacts/{image_asset['storage_relpath']}"
                response = requests.get(
                    f"{base_url}/playbooks/customer-packs/{result['draft_id']}/index.html",
                    timeout=10,
                )
                single_viewer = requests.get(
                    f"{base_url}/api/viewer-document",
                    params={
                        "viewer_path": f"/playbooks/customer-packs/{result['draft_id']}/index.html",
                        "page_mode": "single",
                    },
                    timeout=10,
                )
                multi_viewer = requests.get(
                    f"{base_url}/api/viewer-document",
                    params={
                        "viewer_path": f"/playbooks/customer-packs/{result['draft_id']}/index.html",
                        "page_mode": "multi",
                    },
                    timeout=10,
                )
                asset_response = requests.get(
                    f"{base_url}{asset_url}",
                    timeout=10,
                )

            self.assertEqual(200, response.status_code)
            self.assertIn("CICD 프로세스", response.text)
            self.assertEqual(200, single_viewer.status_code)
            self.assertEqual(200, multi_viewer.status_code)
            single_html = str(single_viewer.json()["html"])
            multi_html = str(multi_viewer.json()["html"])
            self.assertIn('id="slide-001"', single_html)
            self.assertNotIn('id="slide-002"', single_html)
            self.assertIn('class="document-footer-nav"', single_html)
            self.assertIn('href="/playbooks/customer-packs/{}/index.html#slide-002"'.format(result["draft_id"]), single_html)
            self.assertIn(">Quick Nav</summary>", multi_html)
            self.assertIn('id="slide-001"', multi_html)
            self.assertIn('id="slide-002"', multi_html)
            self.assertIn('id="slide-003"', multi_html)
            self.assertIn(asset_url, multi_html)
            self.assertEqual(200, asset_response.status_code)
            self.assertEqual("image/png", asset_response.headers.get("Content-Type"))


if __name__ == "__main__":
    unittest.main()
