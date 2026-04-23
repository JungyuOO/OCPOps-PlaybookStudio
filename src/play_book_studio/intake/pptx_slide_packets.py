from __future__ import annotations

import shutil
from pathlib import Path
from typing import Any

from .models import CustomerPackDraftRecord
from .normalization.builders import (
    _pptx_infer_slide_title,
    _pptx_shape_text,
    _pptx_slide_body_parts,
    _pptx_slide_elements,
)

try:  # pragma: no cover - optional runtime dependency
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception:  # noqa: BLE001
    Presentation = None
    MSO_SHAPE_TYPE = None


CUSTOMER_PACK_SLIDE_PACKET_VERSION = "customer_pack_slide_packet_v1"

_AGENDA_HINTS = ("목차", "agenda", "contents", "table of contents")


def customer_pack_slide_packets_path(books_dir: Path, asset_slug: str) -> Path:
    return books_dir / f"{asset_slug}.slide_packets.json"


def customer_pack_slide_assets_dir(books_dir: Path, asset_slug: str) -> Path:
    return books_dir / f"{asset_slug}.slide-assets"


def build_customer_pack_slide_packets_payload(
    *,
    record: CustomerPackDraftRecord,
    payload: dict[str, Any],
    asset_slug: str,
    book_path: Path,
) -> dict[str, Any]:
    if str(record.request.source_type or "").strip().lower() != "pptx" or Presentation is None:
        return {}
    if str(payload.get("asset_kind") or "").strip() != "customer_pack_manual_book":
        return {}

    capture_path = Path(str(record.capture_artifact_path or "").strip())
    if not capture_path.exists() or not capture_path.is_file():
        return {}

    presentation = Presentation(capture_path)
    books_dir = book_path.parent
    assets_dir = customer_pack_slide_assets_dir(books_dir, asset_slug)
    if assets_dir.exists():
        shutil.rmtree(assets_dir, ignore_errors=True)

    sections = [
        dict(section)
        for section in (payload.get("sections") or [])
        if isinstance(section, dict)
    ]
    matched_section_keys: set[str] = set()
    slides: list[dict[str, Any]] = []
    extracted_assets: list[dict[str, Any]] = []
    total_visual_blocks = 0

    base_viewer_path = (
        str(payload.get("target_viewer_path") or "").strip()
        or f"/playbooks/customer-packs/{record.draft_id}/index.html"
    )
    slide_width = int(getattr(presentation, "slide_width", 0) or 0)
    slide_height = int(getattr(presentation, "slide_height", 0) or 0)

    for slide_index, slide in enumerate(presentation.slides, start=1):
        elements = _pptx_slide_elements(slide)
        title = _pptx_infer_slide_title(slide, slide_index=slide_index, elements=elements)
        body_parts, block_kinds = _pptx_slide_body_parts(elements, title=title)
        body_text = "\n\n".join(part for part in body_parts if str(part).strip()).strip()
        notes_text = _pptx_notes_text(slide)
        matched_section = _match_slide_section(
            sections=sections,
            slide_index=slide_index,
            slide_title=title,
            used_keys=matched_section_keys,
        )
        slide_id = f"{asset_slug}::slide-{slide_index:03d}"
        slide_anchor = f"slide-{slide_index:03d}"
        text_blocks = _slide_text_blocks(elements, slide_id=slide_id, title=title)
        table_blocks = _slide_table_blocks(elements, slide_id=slide_id)
        embedded_assets = _slide_embedded_assets(
            slide=slide,
            slide_index=slide_index,
            slide_id=slide_id,
            asset_slug=asset_slug,
            books_dir=books_dir,
        )
        total_visual_blocks += len(embedded_assets)
        extracted_assets.extend(embedded_assets)
        slide_role = _slide_role(
            title=title,
            body_text=body_text,
            text_blocks=text_blocks,
            table_blocks=table_blocks,
            embedded_assets=embedded_assets,
        )
        slides.append(
            {
                "slide_id": slide_id,
                "slide_anchor": slide_anchor,
                "ordinal": slide_index,
                "title": title,
                "slide_role": slide_role,
                "surface_kind": "slide",
                "source_unit_kind": "slide",
                "origin_method": "native",
                "ocr_status": "not_run",
                "ocr_candidate": bool(embedded_assets),
                "slide_size": {
                    "width": slide_width,
                    "height": slide_height,
                },
                "viewer_path": str((matched_section or {}).get("viewer_path") or base_viewer_path),
                "matched_section_anchor": str((matched_section or {}).get("anchor") or "").strip(),
                "matched_section_key": str((matched_section or {}).get("section_key") or "").strip(),
                "matched_section_heading": str((matched_section or {}).get("heading") or "").strip(),
                "body_text": body_text,
                "notes_text": notes_text,
                "block_kinds": list(dict.fromkeys(str(item).strip() for item in block_kinds if str(item).strip())),
                "text_blocks": text_blocks,
                "table_blocks": table_blocks,
                "embedded_assets": embedded_assets,
                "element_counts": {
                    "text": len(text_blocks),
                    "table": len(table_blocks),
                    "embedded_asset": len(embedded_assets),
                },
            }
        )

    return {
        "artifact_version": CUSTOMER_PACK_SLIDE_PACKET_VERSION,
        "draft_id": str(record.draft_id),
        "asset_slug": asset_slug,
        "book_slug": str(payload.get("book_slug") or asset_slug).strip() or asset_slug,
        "title": str(payload.get("title") or record.plan.title or asset_slug).strip() or asset_slug,
        "surface_kind": "slide_deck",
        "source_unit_kind": "slide",
        "source_type": "pptx",
        "origin_method": "native",
        "ocr_status": "not_run",
        "viewer_path": base_viewer_path,
        "slide_count": len(slides),
        "embedded_asset_count": len(extracted_assets),
        "visual_block_count": total_visual_blocks,
        "ocr_candidate_count": sum(1 for slide in slides if bool(slide.get("ocr_candidate"))),
        "slides": slides,
        "embedded_assets": extracted_assets,
        "slide_size": {
            "width": slide_width,
            "height": slide_height,
        },
    }


def _match_slide_section(
    *,
    sections: list[dict[str, Any]],
    slide_index: int,
    slide_title: str,
    used_keys: set[str],
) -> dict[str, Any] | None:
    normalized_title = _normalize_heading(slide_title)
    ordinal_candidate = sections[slide_index - 1] if 0 < slide_index <= len(sections) else None
    if ordinal_candidate is not None:
        section_key = _section_key(ordinal_candidate)
        if section_key and section_key not in used_keys:
            if normalized_title and _normalize_heading(ordinal_candidate.get("heading")) == normalized_title:
                used_keys.add(section_key)
                return ordinal_candidate
    if not normalized_title:
        return None
    for section in sections:
        section_key = _section_key(section)
        if not section_key or section_key in used_keys:
            continue
        if _normalize_heading(section.get("heading")) != normalized_title:
            continue
        used_keys.add(section_key)
        return section
    return None


def _section_key(section: dict[str, Any]) -> str:
    return str(section.get("section_key") or section.get("anchor") or section.get("ordinal") or "").strip()


def _normalize_heading(value: Any) -> str:
    return " ".join(str(value or "").split()).strip().lower()


def _slide_text_blocks(
    elements: list[dict[str, Any]],
    *,
    slide_id: str,
    title: str,
) -> list[dict[str, Any]]:
    normalized_title = _normalize_heading(title)
    blocks: list[dict[str, Any]] = []
    ordinal = 0
    for element in elements:
        if str(element.get("kind") or "").strip() != "text":
            continue
        ordinal += 1
        text = str(element.get("text") or "").strip()
        if not text:
            continue
        blocks.append(
            {
                "block_id": f"{slide_id}::text-{ordinal:02d}",
                "kind": "text",
                "text": text,
                "is_title_candidate": _normalize_heading(text) == normalized_title,
                "font_size": float(element.get("font_size") or 0.0),
                "bbox": _element_bbox(element),
            }
        )
    return blocks


def _slide_table_blocks(
    elements: list[dict[str, Any]],
    *,
    slide_id: str,
) -> list[dict[str, Any]]:
    blocks: list[dict[str, Any]] = []
    ordinal = 0
    for element in elements:
        if str(element.get("kind") or "").strip() != "table":
            continue
        ordinal += 1
        rows = [
            [str(cell).strip() for cell in row]
            for row in (element.get("rows") or [])
            if isinstance(row, list)
        ]
        row_count = len(rows)
        col_count = max((len(row) for row in rows), default=0)
        blocks.append(
            {
                "block_id": f"{slide_id}::table-{ordinal:02d}",
                "kind": "table",
                "title": str(element.get("title_text") or element.get("text") or "").strip(),
                "body_text": str(element.get("body_text") or "").strip(),
                "row_count": row_count,
                "column_count": col_count,
                "rows": rows,
                "bbox": _element_bbox(element),
            }
        )
    return blocks


def _slide_embedded_assets(
    *,
    slide,
    slide_index: int,
    slide_id: str,
    asset_slug: str,
    books_dir: Path,
) -> list[dict[str, Any]]:
    assets: list[dict[str, Any]] = []
    image_ordinal = 0
    chart_ordinal = 0
    for shape_ordinal, shape in enumerate(slide.shapes, start=1):
        shape_type = getattr(shape, "shape_type", None)
        if MSO_SHAPE_TYPE is not None and shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_ordinal += 1
            try:
                image = shape.image
            except Exception:  # noqa: BLE001
                continue
            extension = str(getattr(image, "ext", "") or "bin").strip().lower() or "bin"
            asset_name = f"slide-{slide_index:03d}-image-{image_ordinal:02d}"
            file_name = f"{asset_name}.{extension}"
            storage_relpath = f"{asset_slug}.slide-assets/{file_name}"
            storage_path = books_dir / storage_relpath
            storage_path.parent.mkdir(parents=True, exist_ok=True)
            storage_path.write_bytes(bytes(getattr(image, "blob", b"") or b""))
            assets.append(
                {
                    "asset_ref": f"{asset_slug}::{asset_name}",
                    "asset_name": asset_name,
                    "asset_kind": "image",
                    "content_type": str(getattr(image, "content_type", "") or ""),
                    "storage_relpath": storage_relpath.replace("\\", "/"),
                    "slide_id": slide_id,
                    "ordinal": shape_ordinal,
                    "bbox": _shape_bbox(shape),
                    "alt": _pptx_shape_text(shape),
                }
            )
            continue
        if bool(getattr(shape, "has_chart", False)):
            chart_ordinal += 1
            assets.append(
                {
                    "asset_ref": f"{asset_slug}::slide-{slide_index:03d}-chart-{chart_ordinal:02d}",
                    "asset_name": f"slide-{slide_index:03d}-chart-{chart_ordinal:02d}",
                    "asset_kind": "chart",
                    "slide_id": slide_id,
                    "ordinal": shape_ordinal,
                    "bbox": _shape_bbox(shape),
                    "alt": _pptx_shape_text(shape),
                }
            )
    return assets


def _element_bbox(element: dict[str, Any]) -> dict[str, int]:
    return {
        "top": int(element.get("top") or 0),
        "left": int(element.get("left") or 0),
        "width": int(element.get("width") or 0),
        "height": int(element.get("height") or 0),
    }


def _shape_bbox(shape) -> dict[str, int]:
    return {
        "top": int(getattr(shape, "top", 0) or 0),
        "left": int(getattr(shape, "left", 0) or 0),
        "width": int(getattr(shape, "width", 0) or 0),
        "height": int(getattr(shape, "height", 0) or 0),
    }


def _pptx_notes_text(slide) -> str:
    try:
        notes_slide = slide.notes_slide
    except Exception:  # noqa: BLE001
        return ""
    lines: list[str] = []
    for shape in getattr(notes_slide, "shapes", []):
        text = _pptx_shape_text(shape)
        if text:
            lines.append(text)
    return "\n\n".join(lines).strip()


def _slide_role(
    *,
    title: str,
    body_text: str,
    text_blocks: list[dict[str, Any]],
    table_blocks: list[dict[str, Any]],
    embedded_assets: list[dict[str, Any]],
) -> str:
    haystack = f"{title}\n{body_text}".lower()
    if any(token in haystack for token in _AGENDA_HINTS):
        return "agenda"
    if embedded_assets and not text_blocks and not table_blocks:
        return "visual_only"
    if table_blocks and len(text_blocks) <= 1:
        return "table_heavy"
    if embedded_assets:
        return "visual_mixed"
    return "content"


__all__ = [
    "CUSTOMER_PACK_SLIDE_PACKET_VERSION",
    "build_customer_pack_slide_packets_payload",
    "customer_pack_slide_assets_dir",
    "customer_pack_slide_packets_path",
]
