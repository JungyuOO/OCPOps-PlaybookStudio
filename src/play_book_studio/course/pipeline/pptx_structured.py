from __future__ import annotations

from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception:  # noqa: BLE001
    Presentation = None
    MSO_SHAPE_TYPE = None


def _shape_text(shape: Any) -> str:
    return str(getattr(shape, "text", "") or "").strip()


def _shape_type_name(shape: Any) -> str:
    if MSO_SHAPE_TYPE is None:
        return "unknown"
    shape_type = getattr(shape, "shape_type", None)
    if shape_type is None:
        return "unknown"
    try:
        return MSO_SHAPE_TYPE(shape_type).name.lower()
    except Exception:  # noqa: BLE001
        return str(shape_type).lower()


def extract_pptx_shapes(pptx_path: Path) -> list[dict[str, Any]]:
    if Presentation is None:
        raise RuntimeError("python-pptx dependency is unavailable for course PPT shape extraction")
    presentation = Presentation(str(pptx_path))
    slide_rows: list[dict[str, Any]] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        shape_rows: list[dict[str, Any]] = []
        text_fragments: list[str] = []
        for shape_index, shape in enumerate(slide.shapes, start=1):
            text = _shape_text(shape)
            if text:
                text_fragments.append(text)
            image_ext = ""
            image_blob = None
            if _shape_type_name(shape) == "picture":
                try:
                    image_ext = str(getattr(shape.image, "ext", "") or "").strip().lower()
                    image_blob = getattr(shape.image, "blob", None)
                except Exception:  # noqa: BLE001
                    image_ext = ""
                    image_blob = None
            shape_rows.append(
                {
                    "shape_index": shape_index,
                    "shape_type": _shape_type_name(shape),
                    "name": str(getattr(shape, "name", "") or "").strip(),
                    "text": text,
                    "left": int(getattr(shape, "left", 0) or 0),
                    "top": int(getattr(shape, "top", 0) or 0),
                    "width": int(getattr(shape, "width", 0) or 0),
                    "height": int(getattr(shape, "height", 0) or 0),
                    "has_table": bool(getattr(shape, "has_table", False)),
                    "image_ext": image_ext,
                    "image_blob": image_blob,
                }
            )
        title = ""
        if slide.shapes.title is not None:
            title = _shape_text(slide.shapes.title)
        slide_rows.append(
            {
                "slide_no": slide_index,
                "title": title or f"Slide {slide_index}",
                "shape_count": len(shape_rows),
                "shapes": shape_rows,
                "text_blob": "\n".join(fragment for fragment in text_fragments if fragment),
            }
        )
    return slide_rows


__all__ = ["extract_pptx_shapes"]
